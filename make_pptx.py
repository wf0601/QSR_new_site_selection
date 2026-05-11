"""Build output/site_selection_deck.pptx — layman-friendly site selection summary."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colours
MCD_RED    = RGBColor(0xDA, 0x29, 0x1C)
MCD_YELLOW = RGBColor(0xFF, 0xC7, 0x2C)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY   = RGBColor(0x75, 0x75, 0x75)
ACCENT     = RGBColor(0x15, 0x65, 0xC0)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def bg(slide, colour: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def box(slide, x, y, w, h, colour: RGBColor, transparency=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = colour
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, colour=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.color.rgb = colour
    return txb


def divider(slide, y, colour=MCD_YELLOW, thickness=Pt(3)):
    line = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.33), thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = colour
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = blank_slide(prs)
    bg(s, DARK)
    box(s, 0, 0, W, Inches(0.18), MCD_YELLOW)
    box(s, 0, H - Inches(0.18), W, Inches(0.18), MCD_RED)

    txt(s, "🍟  McDonald's Tokyo", Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.0),
        size=40, bold=True, colour=MCD_YELLOW, align=PP_ALIGN.LEFT)
    txt(s, "New Site Selection", Inches(0.7), Inches(2.3), Inches(11.9), Inches(0.8),
        size=32, bold=False, colour=WHITE, align=PP_ALIGN.LEFT)
    txt(s, "Using location data + machine learning to find the best spots to open next",
        Inches(0.7), Inches(3.1), Inches(9.5), Inches(0.7),
        size=18, colour=RGBColor(0xBD, 0xBD, 0xBD), align=PP_ALIGN.LEFT)

    txt(s, "Tokyo · 2026", Inches(0.7), H - Inches(0.9), Inches(5), Inches(0.4),
        size=13, colour=MID_GREY, align=PP_ALIGN.LEFT)


def slide_question(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "The Question", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
        size=13, bold=True, colour=MCD_RED, align=PP_ALIGN.LEFT)
    txt(s, "Where should McDonald's open its next restaurant in Tokyo?",
        Inches(0.6), Inches(0.85), Inches(11.8), Inches(1.2),
        size=30, bold=True, colour=DARK, align=PP_ALIGN.LEFT)
    divider(s, Inches(2.1))

    bullets = [
        ("🎯", "Maximise new customers", "Choose locations where demand already exists"),
        ("🚫", "Avoid cannibalisation",  "Stay far enough from existing McDonald's"),
        ("📍", "Pick proven locations",  "Match the profile of spots where McDonald's already succeeds"),
    ]
    for i, (icon, title, desc) in enumerate(bullets):
        y = Inches(2.4) + i * Inches(1.4)
        box(s, Inches(0.6), y, Inches(0.55), Inches(0.55), LIGHT_GREY)
        txt(s, icon, Inches(0.6), y - Inches(0.04), Inches(0.55), Inches(0.6),
            size=22, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(1.3), y, Inches(5), Inches(0.35),
            size=18, bold=True, colour=DARK)
        txt(s, desc,  Inches(1.3), y + Inches(0.35), Inches(10), Inches(0.45),
            size=14, colour=MID_GREY)


def slide_data(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "Step 1 — The Data", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
        size=13, bold=True, colour=MCD_RED)
    txt(s, "All data scraped from Tabelog — Japan's largest restaurant review site",
        Inches(0.6), Inches(0.8), Inches(11.5), Inches(0.6),
        size=22, bold=True, colour=DARK)
    divider(s, Inches(1.55))

    rows = [
        ("McDonald's (own stores)", "275", MCD_RED),
        ("Burger chains (MOS Burger, KFC, Wendy's)", "288", RGBColor(0xE6, 0x51, 0x00)),
        ("Teishoku / beef-bowl chains (Matsuya, Sukiya …)", "960", RGBColor(0x6A, 0x1B, 0x9A)),
        ("Family restaurants (Gusto, Saizeriya)", "281", RGBColor(0x00, 0x69, 0x5C)),
    ]
    for i, (label, count, colour) in enumerate(rows):
        y = Inches(1.85) + i * Inches(1.1)
        box(s, Inches(0.6), y + Inches(0.15), Inches(0.08), Inches(0.55), colour)
        txt(s, label, Inches(0.85), y + Inches(0.15), Inches(8.5), Inches(0.55),
            size=17, colour=DARK)
        txt(s, count + " locations", Inches(9.8), y + Inches(0.15), Inches(2.5), Inches(0.55),
            size=17, bold=True, colour=colour, align=PP_ALIGN.RIGHT)

    txt(s, "Each location includes: GPS coordinates · Tabelog rating · Number of reviews (proxy for popularity)",
        Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
        size=12, colour=MID_GREY)


def slide_weights(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "Step 2 — Not All Competitors Are Equal", Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
        size=13, bold=True, colour=MCD_RED)
    txt(s, "We weight each competitor type by how directly it competes with McDonald's",
        Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        size=22, bold=True, colour=DARK)
    divider(s, Inches(1.55))

    tiers = [
        ("Burger & fast food", "MOS Burger · KFC · Wendy's",
         "Same product, same customer — strongest market signal",
         "1.0", MCD_RED, Inches(6.2)),
        ("Teishoku / beef-bowl", "Matsuya · Sukiya · Yoshinoya · Nakau",
         "Fast-casual overlap — moderate signal",
         "0.5", RGBColor(0x6A, 0x1B, 0x9A), Inches(4.0)),
        ("Family restaurants", "Gusto · Saizeriya",
         "Sit-down dining — weaker signal",
         "0.3", RGBColor(0x00, 0x69, 0x5C), Inches(2.5)),
    ]
    for i, (title, brands, desc, weight, colour, bar_w) in enumerate(tiers):
        y = Inches(1.9) + i * Inches(1.55)
        txt(s, title, Inches(0.6), y, Inches(5), Inches(0.4), size=17, bold=True, colour=colour)
        txt(s, brands, Inches(0.6), y + Inches(0.38), Inches(5.5), Inches(0.35), size=13, colour=MID_GREY)
        txt(s, desc,   Inches(0.6), y + Inches(0.72), Inches(5.5), Inches(0.35), size=13, colour=DARK)
        box(s, Inches(6.5), y + Inches(0.35), bar_w, Inches(0.45), colour)
        txt(s, f"Weight  {weight}", Inches(6.6), y + Inches(0.38), Inches(3.5), Inches(0.4),
            size=15, bold=True, colour=WHITE)

    txt(s, "Popularity score = review count × weight   →   a busy Matsuya counts half as much as a busy MOS Burger",
        Inches(0.6), Inches(6.75), Inches(12), Inches(0.45),
        size=12, colour=MID_GREY)


def slide_methods(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "Step 3 — Two Methods, One Answer", Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
        size=13, bold=True, colour=MCD_RED)
    txt(s, "We combine pattern recognition with gap analysis",
        Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        size=22, bold=True, colour=DARK)
    divider(s, Inches(1.55))

    # Method 1
    box(s, Inches(0.5), Inches(1.75), Inches(5.8), Inches(4.8), LIGHT_GREY)
    txt(s, "Method 1", Inches(0.7), Inches(1.9), Inches(5.5), Inches(0.45),
        size=12, bold=True, colour=MCD_RED)
    txt(s, "Learn from existing stores", Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.7),
        size=20, bold=True, colour=DARK)
    txt(s,
        "We cluster all 275 McDonald's by the competitive landscape around them — "
        "how many rivals are nearby and how popular those rivals are.\n\n"
        "The algorithm finds 5 natural groups of location types.\n\n"
        "A new candidate scores highly if it looks like it belongs to one of those groups.",
        Inches(0.7), Inches(3.05), Inches(5.4), Inches(3.2),
        size=13, colour=DARK)
    txt(s, "→  Cluster Affinity Score (50%)", Inches(0.7), Inches(6.15), Inches(5.4), Inches(0.35),
        size=13, bold=True, colour=MCD_RED)

    # Method 2
    box(s, Inches(7.0), Inches(1.75), Inches(5.8), Inches(4.8), LIGHT_GREY)
    txt(s, "Method 2", Inches(7.2), Inches(1.9), Inches(5.5), Inches(0.45),
        size=12, bold=True, colour=ACCENT)
    txt(s, "Find the gaps", Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.7),
        size=20, bold=True, colour=DARK)
    txt(s,
        "We look for areas where food competition is active "
        "(proving people eat out there) but no McDonald's exists nearby.\n\n"
        "High competitor popularity + far from any existing McDonald's "
        "= underserved pocket of demand.\n\n"
        "This is the infill opportunity.",
        Inches(7.2), Inches(3.05), Inches(5.4), Inches(3.2),
        size=13, colour=DARK)
    txt(s, "→  Infill Score (30%)", Inches(7.2), Inches(6.15), Inches(5.4), Inches(0.35),
        size=13, bold=True, colour=ACCENT)

    txt(s, "+  Distance Buffer (20%) — rewards safe spacing from existing stores",
        Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
        size=12, colour=MID_GREY)


def slide_guardrails(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "Step 4 — Safety Rules", Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
        size=13, bold=True, colour=MCD_RED)
    txt(s, "Before scoring, every candidate must pass three filters",
        Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        size=22, bold=True, colour=DARK)
    divider(s, Inches(1.55))

    rules = [
        ("🚫", "No cannibalisation",
         "Candidate must be at least 0.8 km from the nearest existing McDonald's",
         "If we open too close to ourselves, we just steal our own customers"),
        ("📡", "Market must exist",
         "At least 1 competitor must be within 1 km; nearest competitor within 3 km",
         "No competition = no proven demand = not worth opening there"),
        ("📊", "Real activity required",
         "Weighted competitive demand score must be greater than zero",
         "Avoids ghost locations with zero nearby food footfall"),
    ]
    for i, (icon, title, rule, reason) in enumerate(rules):
        y = Inches(1.85) + i * Inches(1.55)
        txt(s, icon, Inches(0.5), y, Inches(0.65), Inches(0.7), size=28, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(1.3), y, Inches(10), Inches(0.45), size=18, bold=True, colour=DARK)
        txt(s, rule,  Inches(1.3), y + Inches(0.42), Inches(10.5), Inches(0.4), size=14, colour=DARK)
        txt(s, f"Why: {reason}", Inches(1.3), y + Inches(0.82), Inches(10.5), Inches(0.45),
            size=12, colour=MID_GREY)

    txt(s, f"~8,500 grid points across Tokyo evaluated — filtered down to the viable candidates before scoring",
        Inches(0.6), Inches(6.8), Inches(12), Inches(0.45), size=12, colour=MID_GREY)


def slide_results(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "Step 5 — Results", Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
        size=13, bold=True, colour=MCD_RED)
    txt(s, "Top 5 candidate sites — Meguro / Setagaya area dominates",
        Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        size=22, bold=True, colour=DARK)
    divider(s, Inches(1.55))

    headers = ["Rank", "Area (approx.)", "Score", "Cluster fit", "Infill demand", "Gap from nearest McDonald's"]
    col_x   = [Inches(0.5), Inches(1.2), Inches(4.8), Inches(6.0), Inches(7.8), Inches(9.6)]
    col_w   = [Inches(0.6), Inches(3.5), Inches(1.1), Inches(1.7), Inches(1.7), Inches(3.1)]

    # Header row
    box(s, Inches(0.45), Inches(1.75), Inches(12.4), Inches(0.45), DARK)
    for hdr, x, w in zip(headers, col_x, col_w):
        txt(s, hdr, x, Inches(1.78), w, Inches(0.4), size=11, bold=True, colour=WHITE)

    rows = [
        ("#1", "Meguro-ku (Yutenji area)", "0.889", "100%", "80%", "1.48 km"),
        ("#2", "Meguro-ku (Yutenji area)", "0.872", "95%",  "78%", "1.62 km"),
        ("#3", "Meguro-ku (Yutenji area)", "0.863", "100%", "77%", "1.34 km"),
        ("#4", "Setagaya-ku",              "0.849", "87%",  "100%","1.14 km"),
        ("#5", "Meguro-ku",               "0.842", "100%", "83%", "0.94 km"),
    ]
    for i, row in enumerate(rows):
        y = Inches(2.25) + i * Inches(0.82)
        fill = LIGHT_GREY if i % 2 == 0 else WHITE
        box(s, Inches(0.45), y, Inches(12.4), Inches(0.78), fill)
        if i == 0:
            box(s, Inches(0.45), y, Inches(0.07), Inches(0.78), MCD_RED)
        for val, x, w in zip(row, col_x, col_w):
            txt(s, val, x, y + Inches(0.2), w, Inches(0.45),
                size=13, bold=(i == 0), colour=DARK)

    txt(s, "Full ranked list of 20 sites saved to output/top_candidates.csv",
        Inches(0.6), Inches(6.85), Inches(12), Inches(0.4), size=12, colour=MID_GREY)


def slide_outputs(prs):
    s = blank_slide(prs)
    bg(s, WHITE)
    box(s, 0, 0, W, Inches(0.12), MCD_RED)

    txt(s, "What You Get", Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
        size=13, bold=True, colour=MCD_RED)
    txt(s, "Two ready-to-use deliverables",
        Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        size=22, bold=True, colour=DARK)
    divider(s, Inches(1.55))

    # Output 1
    box(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5), LIGHT_GREY)
    txt(s, "📄", Inches(0.8), Inches(2.0), Inches(1), Inches(0.9), size=36, align=PP_ALIGN.CENTER)
    txt(s, "top_candidates.csv", Inches(1.7), Inches(2.05), Inches(4.4), Inches(0.5),
        size=17, bold=True, colour=DARK)
    txt(s,
        "A spreadsheet with the top 20 candidate sites.\n\n"
        "Each row shows the location coordinates, final score, "
        "and a breakdown of all three score components.\n\n"
        "Open in Excel or Google Sheets.",
        Inches(0.7), Inches(2.7), Inches(5.4), Inches(3.0),
        size=13, colour=DARK)

    # Output 2
    box(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.5), LIGHT_GREY)
    txt(s, "🗺️", Inches(7.3), Inches(2.0), Inches(1), Inches(0.9), size=36, align=PP_ALIGN.CENTER)
    txt(s, "interactive_map.html", Inches(8.2), Inches(2.05), Inches(4.4), Inches(0.5),
        size=17, bold=True, colour=DARK)
    txt(s,
        "An interactive map that opens in any web browser.\n\n"
        "Toggle layers on/off to see McDonald's locations, "
        "each competitor tier (burger / teishoku / family), "
        "and the top 20 candidate sites with detailed score pop-ups.",
        Inches(7.2), Inches(2.7), Inches(5.4), Inches(3.0),
        size=13, colour=DARK)

    txt(s, "Run:  python3 site_selection.py   to regenerate both files with any updated data",
        Inches(0.6), Inches(6.85), Inches(12), Inches(0.4), size=12, colour=MID_GREY)


def slide_summary(prs):
    s = blank_slide(prs)
    bg(s, DARK)
    box(s, 0, 0, W, Inches(0.18), MCD_YELLOW)
    box(s, 0, H - Inches(0.18), W, Inches(0.18), MCD_RED)

    txt(s, "In Plain English", Inches(0.7), Inches(0.55), Inches(12), Inches(0.6),
        size=14, bold=True, colour=MCD_YELLOW, align=PP_ALIGN.LEFT)
    txt(s, "We asked the data:\n\"Where is there food demand in Tokyo that McDonald's isn't serving?\"",
        Inches(0.7), Inches(1.2), Inches(11.9), Inches(1.6),
        size=26, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)

    points = [
        "Scraped 1,800+ restaurant locations from Tabelog",
        "Taught an algorithm what a good McDonald's location looks like",
        "Weighted competitors by how directly they compete",
        "Scanned 8,500 points across Tokyo to find the best gaps",
        "Delivered 20 ranked sites with an interactive map",
    ]
    for i, p in enumerate(points):
        y = Inches(3.0) + i * Inches(0.72)
        box(s, Inches(0.7), y + Inches(0.12), Inches(0.08), Inches(0.42), MCD_YELLOW)
        txt(s, p, Inches(1.0), y, Inches(11.5), Inches(0.65),
            size=16, colour=RGBColor(0xE0, 0xE0, 0xE0))

    txt(s, "Top recommendation: Meguro / Setagaya area  ·  Score 0.889",
        Inches(0.7), Inches(6.8), Inches(12), Inches(0.45),
        size=14, bold=True, colour=MCD_YELLOW)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_question(prs)
    slide_data(prs)
    slide_weights(prs)
    slide_methods(prs)
    slide_guardrails(prs)
    slide_results(prs)
    slide_outputs(prs)
    slide_summary(prs)

    out = "output/site_selection_deck.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
